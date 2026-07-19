# app/utils/file_extractor.py
# ---------------------------------------------------------------------------
# Upgraded extractor for AiRev / Agent@4
#
# Supported student-solution formats:
#   1. Direct typing (handled in route, not here)
#   2. PDF — text-based       -> pypdf
#   3. PDF — scanned/photo    -> auto-rasterize -> Claude vision OCR
#   4. DOCX                   -> python-docx
#   5. DOC  (legacy)          -> friendly "save as .docx" message
#   6. TXT / MD               -> UTF-8 decode
#   7. RTF                    -> control-word strip
#   8. Images (handwritten/photographed notes):
#         JPG, JPEG, PNG, WEBP, HEIC, HEIF -> Claude vision OCR
#
# Cost guards:
#   - MAX_OCR_PAGES env (default 5) caps PDF rasterization
#   - MAX_FILE_BYTES env (default 10 MB) rejects oversized uploads
#   - HEIC support is conditional on pillow-heif being installed
# ---------------------------------------------------------------------------

import io
import os
import re
import base64
import logging
from typing import Tuple, List

import httpx

logger = logging.getLogger(__name__)

# ---------- config ---------------------------------------------------------

MAX_FILE_BYTES = int(os.getenv("MAX_FILE_BYTES", str(10 * 1024 * 1024)))   # 10 MB
MAX_OCR_PAGES = int(os.getenv("MAX_OCR_PAGES", "5"))                       # OCR cap
OCR_MODEL = os.getenv("OCR_MODEL", "claude-sonnet-4-5")                    # vision-capable

IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".webp", ".heic", ".heif"}
TEXT_EXTS  = {".txt", ".md"}
CODE_EXTS  = {".py", ".sql", ".js", ".ts", ".jsx", ".tsx", ".java", ".c",
              ".cpp", ".r", ".json", ".html", ".css", ".sh", ".yaml", ".yml"}
SHEET_MAX_ROWS  = 300     # per sheet — enough for any coursework workbook
SHEET_MAX_COLS  = 40
SHEET_MAX_CHARS = 60000   # whole-workbook render cap; truncation is stated
ZIP_MAX_FILES   = 25
ZIP_MAX_TOTAL   = 60 * 1024 * 1024   # unpacked-bytes bomb guard


# ---------- public entry ---------------------------------------------------

def extract_text_from_url(file_url: str, file_name: str = "") -> Tuple[str, str]:
    """
    Returns (extracted_text, reason).
    On success: ('extracted text...', '')
    On failure: ('', 'short reason for log')
    """
    if not file_url:
        return "", "no file_url provided"

    data, why = _download_file(file_url)
    if data is None:
        return "", why

    if len(data) > MAX_FILE_BYTES:
        return "", f"file too large ({len(data) // 1024} KB > {MAX_FILE_BYTES // 1024} KB)"

    name = (file_name or file_url).lower()
    ext = "." + name.rsplit(".", 1)[-1] if "." in name else ""

    try:
        # PDFs ------------------------------------------------------------
        if ext == ".pdf":
            text, why = _extract_pdf(data)
            if text:
                return text, ""
            # Empty -> probably scanned. Try OCR.
            logger.info("PDF text empty (%s) -> falling back to vision OCR", why)
            return _extract_scanned_pdf(data)

        # Word ------------------------------------------------------------
        if ext == ".docx":
            return _extract_docx(data)
        if ext == ".doc":
            return "", (
                "Legacy .doc format isn't supported. "
                "Please open the file in Word, choose 'Save As', "
                "select 'Word Document (.docx)', and re-upload."
            )

        # Plain text ------------------------------------------------------
        if ext in TEXT_EXTS:
            return _clean(data.decode("utf-8", errors="ignore")), ""
        if ext == ".rtf":
            return _extract_rtf(data)

        # Spreadsheets (CASA-class assignments demand Excel workbooks) -----
        if ext in {".xlsx", ".xlsm", ".xltx"}:
            return _extract_xlsx(data)
        if ext in {".csv", ".tsv"}:
            return _extract_csv(data, ext)

        # Presentations ----------------------------------------------------
        if ext in {".pptx", ".potx"}:
            return _extract_pptx(data)

        # Archives (capstone UI promises ZIP — honor it) -------------------
        if ext == ".zip":
            return _extract_zip(data)

        # Code / notebooks (capstones: "build a decisioning engine") -------
        if ext == ".ipynb":
            return _extract_ipynb(data)
        if ext in CODE_EXTS:
            body = _clean(data.decode("utf-8", errors="ignore"))[:SHEET_MAX_CHARS]
            return (f"[Code file: {name.rsplit('/', 1)[-1]}]\n{body}", "") if body \
                else ("", "code file was empty")

        # Images (handwritten notes) -------------------------------------
        if ext in IMAGE_EXTS:
            return _extract_image(data, ext)

        # Unknown ext -> sniff. Try PDF, DOCX, then image, then bytes-as-text.
        for fn in (_extract_pdf, _extract_docx):
            text, _ = fn(data)
            if text:
                return text, ""
        if _looks_like_image(data):
            return _extract_image(data, ".png")
        return _clean(data.decode("utf-8", errors="ignore")), ""

    except Exception as e:
        logger.exception("extraction crashed")
        return "", f"extraction failed: {type(e).__name__}"


# ---------- download (kept compatible with original) -----------------------

def _download_file(file_url: str) -> Tuple[bytes, str]:
    """Returns (bytes, reason). bytes is None on failure."""
    cloud_name = os.getenv("CLOUDINARY_CLOUD_NAME")
    api_key    = os.getenv("CLOUDINARY_API_KEY")
    api_secret = os.getenv("CLOUDINARY_API_SECRET")

    try:
        with httpx.Client(timeout=30, follow_redirects=True) as client:
            r = client.get(file_url)
            if r.status_code == 200:
                return r.content, ""
            # Authenticated Cloudinary fallback
            if r.status_code in (401, 403) and cloud_name and api_key and api_secret:
                m = re.search(
                    rf"https://res\.cloudinary\.com/{re.escape(cloud_name)}/[^/]+/authenticated/(.+)",
                    file_url,
                )
                if m:
                    public_id = m.group(1)
                    try:
                        import cloudinary, cloudinary.utils
                        cloudinary.config(
                            cloud_name=cloud_name,
                            api_key=api_key,
                            api_secret=api_secret,
                        )
                        signed_url, _ = cloudinary.utils.cloudinary_url(
                            public_id, type="authenticated", sign_url=True,
                        )
                        r2 = client.get(signed_url)
                        if r2.status_code == 200:
                            return r2.content, ""
                        return None, f"signed download HTTP {r2.status_code}"
                    except Exception as e:
                        return None, f"cloudinary sign failed: {e}"
                return None, f"download HTTP {r.status_code} (no Cloudinary credentials to retry)"
            return None, f"download HTTP {r.status_code}"
    except Exception as e:
        return None, f"download failed: {e}"


# ---------- PDF ------------------------------------------------------------

def _extract_pdf(data: bytes) -> Tuple[str, str]:
    try:
        from pypdf import PdfReader
    except ImportError:
        return "", "pypdf not installed"
    try:
        reader = PdfReader(io.BytesIO(data))
        chunks = []
        for page in reader.pages:
            try:
                chunks.append(page.extract_text() or "")
            except Exception:
                continue
        text = _clean("\n".join(chunks))
        if not text:
            return "", "PDF parsed but no extractable text (likely scanned image)"
        return text, ""
    except Exception as e:
        return "", f"pdf parse error: {e}"


def _extract_scanned_pdf(data: bytes) -> Tuple[str, str]:
    """Rasterize first N pages and OCR via Claude vision."""
    try:
        import pypdfium2 as pdfium
    except ImportError:
        return "", "scanned PDF detected — install pypdfium2 to enable OCR"

    try:
        pdf = pdfium.PdfDocument(data)
    except Exception as e:
        return "", f"pdf rasterize open failed: {e}"

    n = min(len(pdf), MAX_OCR_PAGES)
    if n == 0:
        return "", "PDF has zero pages"

    images_b64: List[Tuple[str, str]] = []  # (media_type, base64)
    try:
        for i in range(n):
            page = pdf[i]
            # 144 DPI is a good handwriting/print balance
            bitmap = page.render(scale=2.0)
            pil_img = bitmap.to_pil()
            buf = io.BytesIO()
            pil_img.save(buf, format="PNG", optimize=True)
            images_b64.append(("image/png", base64.b64encode(buf.getvalue()).decode()))
    except Exception as e:
        return "", f"pdf rasterize page failed: {e}"
    finally:
        try:
            pdf.close()
        except Exception:
            pass

    text, why = _ocr_with_claude(images_b64, kind="scanned PDF")
    if not text:
        return "", why
    if len(pdf) > MAX_OCR_PAGES:
        text += (
            f"\n\n[Note: only the first {MAX_OCR_PAGES} pages were OCR-processed "
            f"out of {len(pdf)} total. Increase MAX_OCR_PAGES to read more.]"
        )
    return text, ""


# ---------- DOCX -----------------------------------------------------------

def _extract_docx(data: bytes) -> Tuple[str, str]:
    try:
        from docx import Document
    except ImportError:
        return "", "python-docx not installed"
    try:
        doc = Document(io.BytesIO(data))
        parts = [p.text for p in doc.paragraphs if p.text]
        # Also pull table cells (case studies often use them)
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text:
                        parts.append(cell.text)
        text = _clean("\n".join(parts))
        if not text:
            return "", "DOCX parsed but empty"
        return text, ""
    except Exception as e:
        return "", f"docx parse error: {e}"


# ---------- RTF ------------------------------------------------------------

def _extract_rtf(data: bytes) -> Tuple[str, str]:
    """Strip RTF control words. Good enough for plain answer text."""
    try:
        raw = data.decode("utf-8", errors="ignore")
        # Drop binary blocks
        raw = re.sub(r"\\pict[^}]*\}", "", raw)
        # Drop control words like \rtf1, \par, \fs24, \'e9
        raw = re.sub(r"\\[a-zA-Z]+-?\d*\s?", " ", raw)
        raw = re.sub(r"\\'[0-9a-fA-F]{2}", "", raw)
        # Drop braces
        raw = re.sub(r"[{}]", "", raw)
        text = _clean(raw)
        if not text:
            return "", "RTF parsed but empty"
        return text, ""
    except Exception as e:
        return "", f"rtf parse error: {e}"


# ---------- Image OCR ------------------------------------------------------

def _extract_image(data: bytes, ext: str) -> Tuple[str, str]:
    """OCR a single image (handwritten or printed)."""
    media_type = _media_type_for_ext(ext)

    # Convert HEIC/HEIF -> PNG so Claude can read it
    if ext in {".heic", ".heif"}:
        try:
            import pillow_heif
            from PIL import Image
            pillow_heif.register_heif_opener()
            img = Image.open(io.BytesIO(data))
            buf = io.BytesIO()
            img.convert("RGB").save(buf, format="PNG", optimize=True)
            data = buf.getvalue()
            media_type = "image/png"
        except ImportError:
            return "", (
                "HEIC images need pillow-heif. Either install it, "
                "or convert the photo to JPG/PNG and re-upload."
            )
        except Exception as e:
            return "", f"heic conversion failed: {e}"

    b64 = base64.b64encode(data).decode()
    return _ocr_with_claude([(media_type, b64)], kind="photographed notes")


def _ocr_with_claude(images: List[Tuple[str, str]], kind: str) -> Tuple[str, str]:
    """Run vision OCR on one or more images. Returns (text, reason)."""
    api_key = os.getenv("ANTHROPIC_API_KEY")
    if not api_key:
        return "", "ANTHROPIC_API_KEY not set — cannot OCR"

    try:
        import anthropic
    except ImportError:
        return "", "anthropic SDK not installed"

    content = []
    for media_type, b64 in images:
        content.append({
            "type": "image",
            "source": {"type": "base64", "media_type": media_type, "data": b64},
        })
    content.append({
        "type": "text",
        "text": (
            f"The image(s) above are a student's {kind} for a case-study answer. "
            "Transcribe ALL handwritten or printed text you can read, exactly as written. "
            "Preserve paragraph breaks and bullet points. Do not summarize, do not add "
            "commentary, do not correct grammar. If multiple pages are shown, separate them "
            "with a blank line. If a section is unreadable, write [unreadable] in its place. "
            "Output only the transcription — no preamble."
        ),
    })

    try:
        client = anthropic.Anthropic(api_key=api_key)
        msg = client.messages.create(
            model=OCR_MODEL,
            max_tokens=4000,
            messages=[{"role": "user", "content": content}],
        )
        text_parts = [b.text for b in msg.content if getattr(b, "type", "") == "text"]
        text = _clean("\n".join(text_parts))
        if not text:
            return "", "OCR returned empty text"
        return text, ""
    except Exception as e:
        logger.exception("vision OCR failed")
        return "", f"OCR failed: {type(e).__name__}"


# ---------- Spreadsheets ----------------------------------------------------

def _extract_xlsx(data: bytes) -> Tuple[str, str]:
    """Render a workbook as reviewable text: values AND formulas, sheet by
    sheet. Formulas matter — an assignment asking for computed ratios must be
    judged on whether the student actually computed them, not typed them."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        return "", "openpyxl not installed"
    try:
        wb_formulas = load_workbook(io.BytesIO(data), data_only=False, read_only=True)
        wb_values   = load_workbook(io.BytesIO(data), data_only=True,  read_only=True)
    except Exception as e:
        return "", f"xlsx parse error: {e}"

    out, truncated = [], False
    for sheet_name in wb_formulas.sheetnames:
        wsf, wsv = wb_formulas[sheet_name], wb_values[sheet_name]
        out.append(f"=== SHEET: {sheet_name} ===")
        for r_idx, (row_f, row_v) in enumerate(zip(
                wsf.iter_rows(max_row=SHEET_MAX_ROWS, max_col=SHEET_MAX_COLS),
                wsv.iter_rows(max_row=SHEET_MAX_ROWS, max_col=SHEET_MAX_COLS))):
            cells = []
            for cf, cv in zip(row_f, row_v):
                if cf.value is None and cv.value is None:
                    cells.append("")
                    continue
                formula = str(cf.value) if isinstance(cf.value, str) and str(cf.value).startswith("=") else None
                value = cv.value if cv.value is not None else cf.value
                cells.append(f"{value} [{formula}]" if formula else str(value))
            line = " | ".join(cells).rstrip(" |")
            if line.strip():
                out.append(line)
        if (wsf.max_row or 0) > SHEET_MAX_ROWS or (wsf.max_column or 0) > SHEET_MAX_COLS:
            truncated = True
        if sum(len(x) for x in out) > SHEET_MAX_CHARS:
            truncated = True
            break

    text = _clean("\n".join(out))[:SHEET_MAX_CHARS]
    if truncated:
        text += "\n\n[Note: workbook truncated for review — very large sheets are rendered partially.]"
    if not text or text.startswith("=== SHEET") and len(text) < 40:
        return "", "workbook parsed but contained no data"
    return text, ""


def _extract_csv(data: bytes, ext: str) -> Tuple[str, str]:
    raw = data.decode("utf-8-sig", errors="ignore")
    lines = raw.splitlines()[:SHEET_MAX_ROWS]
    text = _clean("\n".join(lines))[:SHEET_MAX_CHARS]
    if len(raw.splitlines()) > SHEET_MAX_ROWS:
        text += f"\n\n[Note: showing first {SHEET_MAX_ROWS} rows of {len(raw.splitlines())}.]"
    return (text, "") if text else ("", "csv was empty")


# ---------- Presentations ---------------------------------------------------

def _extract_pptx(data: bytes) -> Tuple[str, str]:
    try:
        from pptx import Presentation
    except ImportError:
        return "", "python-pptx not installed"
    try:
        prs = Presentation(io.BytesIO(data))
        out = []
        for i, slide in enumerate(prs.slides, 1):
            parts = [f"--- Slide {i} ---"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = "".join(run.text for run in para.runs).strip()
                        if t:
                            parts.append(t)
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    parts.append(f"[Speaker notes] {notes}")
            out.append("\n".join(parts))
        text = _clean("\n\n".join(out))[:SHEET_MAX_CHARS]
        return (text, "") if text else ("", "pptx parsed but empty")
    except Exception as e:
        return "", f"pptx parse error: {e}"


# ---------- Archives ---------------------------------------------------------

def _extract_zip(data: bytes) -> Tuple[str, str]:
    """Unpack in memory and extract every supported file inside. The capstone
    UI promises ZIP support — this honors it. Bomb-guarded."""
    import zipfile
    try:
        zf = zipfile.ZipFile(io.BytesIO(data))
    except Exception as e:
        return "", f"zip open failed: {e}"

    infos = [i for i in zf.infolist()
             if not i.is_dir() and not i.filename.startswith("__MACOSX")]
    if not infos:
        return "", "zip was empty"
    if sum(i.file_size for i in infos) > ZIP_MAX_TOTAL:
        return "", f"zip unpacks beyond {ZIP_MAX_TOTAL // (1024*1024)}MB — too large to review"

    out, skipped = [], []
    for info in infos[:ZIP_MAX_FILES]:
        inner_name = info.filename
        inner_ext = "." + inner_name.rsplit(".", 1)[-1].lower() if "." in inner_name else ""
        if inner_ext == ".zip":
            skipped.append(inner_name)      # no nested archives
            continue
        try:
            inner = zf.read(info)
        except Exception:
            skipped.append(inner_name)
            continue
        text, why = _extract_inner(inner, inner_name, inner_ext)
        if text:
            out.append(f"===== FILE: {inner_name} =====\n{text}")
        else:
            skipped.append(f"{inner_name} ({why})" if why else inner_name)

    if len(infos) > ZIP_MAX_FILES:
        skipped.append(f"... and {len(infos) - ZIP_MAX_FILES} more files beyond the {ZIP_MAX_FILES}-file cap")
    combined = _clean("\n\n".join(out))
    if not combined:
        return "", ("zip contained no readable files; skipped: " + ", ".join(skipped[:8]))
    if skipped:
        combined += "\n\n[Files in the zip that could not be read: " + ", ".join(skipped[:10]) + "]"
    return combined, ""


def _extract_inner(data: bytes, name: str, ext: str) -> Tuple[str, str]:
    """Extract one file from inside a zip using the standard handlers."""
    if ext == ".pdf":
        text, why = _extract_pdf(data)
        return (text, why) if text else ("", why)   # no OCR inside zips — cost guard
    if ext == ".docx":
        return _extract_docx(data)
    if ext in {".xlsx", ".xlsm", ".xltx"}:
        return _extract_xlsx(data)
    if ext in {".csv", ".tsv"}:
        return _extract_csv(data, ext)
    if ext in {".pptx", ".potx"}:
        return _extract_pptx(data)
    if ext == ".ipynb":
        return _extract_ipynb(data)
    if ext in TEXT_EXTS or ext in CODE_EXTS:
        body = _clean(data.decode("utf-8", errors="ignore"))[:SHEET_MAX_CHARS]
        return (body, "") if body else ("", "empty")
    return "", f"unsupported inside zip ({ext or 'no extension'})"


# ---------- Notebooks --------------------------------------------------------

def _extract_ipynb(data: bytes) -> Tuple[str, str]:
    import json as _json
    try:
        nb = _json.loads(data.decode("utf-8", errors="ignore"))
        out = []
        for cell in nb.get("cells", []):
            kind = cell.get("cell_type")
            src = "".join(cell.get("source", [])).strip()
            if not src:
                continue
            out.append(f"[{kind} cell]\n{src}" if kind else src)
        text = _clean("\n\n".join(out))[:SHEET_MAX_CHARS]
        return (text, "") if text else ("", "notebook had no content cells")
    except Exception as e:
        return "", f"ipynb parse error: {e}"


# ---------- helpers --------------------------------------------------------

def _media_type_for_ext(ext: str) -> str:
    return {
        ".jpg":  "image/jpeg",
        ".jpeg": "image/jpeg",
        ".png":  "image/png",
        ".webp": "image/webp",
        ".heic": "image/heic",
        ".heif": "image/heif",
    }.get(ext, "image/png")


def _looks_like_image(data: bytes) -> bool:
    if len(data) < 12:
        return False
    return (
        data[:3] == b"\xff\xd8\xff"            # JPEG
        or data[:8] == b"\x89PNG\r\n\x1a\n"    # PNG
        or data[:4] == b"RIFF" and data[8:12] == b"WEBP"  # WEBP
    )


def _clean(text: str) -> str:
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()